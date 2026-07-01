#!/usr/bin/env python3
"""
Generate a branded PowerPoint status report from AI-generated narrative content.
Reads JSON data from stdin, outputs PPTX to the path specified as first argument.
"""

import sys
import json
import os
import re
import copy
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

FONT_DIR = os.path.join(os.path.dirname(__file__), '..', '..', 'client', 'public', 'fonts')
FONT_NAME = 'Avenir Next LT Pro'

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

def set_font(run, size=10, bold=False, color=None, name=FONT_NAME, italic=False):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color if isinstance(color, RGBColor) else hex_to_rgb(color)

def set_cell_text(cell, text, size=9, bold=False, color=None, bg_color=None, alignment=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    set_font(run, size=size, bold=bold, color=color)
    cell.vertical_anchor = valign
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color if isinstance(bg_color, RGBColor) else hex_to_rgb(bg_color)

def add_accent_bar(slide, primary_color, left=0, top=0, width=None, height=Inches(0.08)):
    if width is None:
        width = SLIDE_WIDTH
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(primary_color)
    shape.line.fill.background()
    return shape

SECTION_ALIASES = {
    'Key Accomplishments': [
        'key accomplishments', 'accomplishments', 'progress & accomplishments',
        'progress and accomplishments', 'work completed', 'completed work',
        'key highlights', 'highlights', 'achievements',
    ],
    'Progress Summary': [
        'progress summary', 'executive summary', 'summary', 'overview',
        'project summary', 'status summary',
    ],
    'Upcoming Activities': [
        'upcoming activities', 'next steps', 'upcoming', 'looking ahead',
        'next period', 'planned activities', 'upcoming work',
    ],
    'Risks, Issues & Key Decisions (RAIDD)': [
        'risks, issues & key decisions', 'raidd', 'risks and issues',
        'risks, issues & key decisions (raidd)', 'risk summary',
        'raidd summary', 'raidd log',
    ],
}

def _normalize_section_name(raw_name):
    """Map a raw AI section header to a canonical section name."""
    lower = raw_name.lower().strip()
    lower = re.sub(r'\s*\(.*?\)\s*$', '', lower).strip()
    for canonical, aliases in SECTION_ALIASES.items():
        if lower in aliases:
            return canonical
        for alias in aliases:
            if alias in lower or lower in alias:
                return canonical
    return raw_name

def parse_markdown_sections(md_text):
    """Parse AI-generated markdown into structured sections with fuzzy header matching."""
    sections = {}
    current_section = None
    current_content = []

    for line in md_text.split('\n'):
        header_match = re.match(r'^(#{1,3})\s+(.+)$', line)
        if header_match:
            if current_section:
                sections[current_section] = '\n'.join(current_content).strip()
            raw_name = header_match.group(2).strip()
            current_section = _normalize_section_name(raw_name)
            current_content = []
        else:
            current_content.append(line)

    if current_section:
        sections[current_section] = '\n'.join(current_content).strip()

    print(f"[PPTX] Parsed AI sections: {list(sections.keys())}", file=sys.stderr)
    return sections

def parse_bullet_items(text):
    """Parse markdown bullet items into structured data with bold titles and descriptions."""
    items = []
    current_item = None

    for line in text.split('\n'):
        stripped = line.strip()
        if not stripped:
            continue

        if stripped.startswith('- ') or stripped.startswith('* ') or stripped.startswith('• '):
            if current_item:
                items.append(current_item)
            content = stripped[2:].strip()
            bold_match = re.match(r'\*\*(.+?)\*\*\s*[-–—:]?\s*(.*)', content)
            if bold_match:
                current_item = {
                    'title': bold_match.group(1).strip(),
                    'description': bold_match.group(2).strip(),
                    'sub_items': []
                }
            else:
                current_item = {
                    'title': content,
                    'description': '',
                    'sub_items': []
                }
        elif stripped.startswith('  - ') or stripped.startswith('  * ') or stripped.startswith('  • '):
            if current_item:
                sub_content = stripped.lstrip(' -•*').strip()
                current_item['sub_items'].append(sub_content)
        elif current_item:
            if current_item['description']:
                current_item['description'] += ' ' + stripped
            else:
                current_item['description'] = stripped

    if current_item:
        items.append(current_item)

    return items

def render_markdown_text(tf, text, primary_color, size=10, start_fresh=True):
    """Render markdown text with bold formatting into a text frame."""
    if start_fresh and tf.paragraphs[0].text == '':
        started = False
    else:
        started = True

    lines = text.split('\n')
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue

        if started:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
            started = True

        p.space_after = Pt(3)
        p.space_before = Pt(2)

        if stripped.startswith('- ') or stripped.startswith('• ') or stripped.startswith('* '):
            content = stripped[2:].strip()
            render_inline_bold(p, f"• {content}", size=size, primary_color=primary_color)
        elif stripped.startswith('  - ') or stripped.startswith('  • '):
            content = stripped.lstrip(' -•*').strip()
            p.level = 1
            render_inline_bold(p, f"  – {content}", size=size - 1, primary_color=primary_color)
        else:
            render_inline_bold(p, stripped, size=size, primary_color=primary_color)

def render_inline_bold(p, text, size=10, primary_color=None):
    """Render text with **bold** markdown formatting into a paragraph."""
    parts = re.split(r'(\*\*[^*]+\*\*)', text)
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            run = p.add_run()
            run.text = part[2:-2]
            set_font(run, size=size, bold=True, color=primary_color)
        else:
            if part:
                run = p.add_run()
                run.text = part
                set_font(run, size=size)

def create_title_slide(prs, data, primary_color, secondary_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = hex_to_rgb(primary_color)
    bg_shape.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.0), SLIDE_WIDTH, Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = hex_to_rgb(secondary_color)
    accent.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = data.get('projectName', 'Project Status Report')
    set_font(run, size=36, bold=True, color=RGBColor(255, 255, 255))

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = "STATUS REPORT"
    set_font(run2, size=20, color=RGBColor(220, 220, 220))

    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1.2))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    run = p.add_run()
    run.text = data.get('clientName', '')
    set_font(run, size=16, bold=True, color=RGBColor(255, 255, 255))

    period_start = data.get('periodStart', '')
    period_end = data.get('periodEnd', '')
    report_date = data.get('reportDate', datetime.now().strftime('%B %d, %Y'))
    period_text = f"Period: {period_start} to {period_end}" if period_start and period_end else report_date

    p2 = tf2.add_paragraph()
    run2 = p2.add_run()
    run2.text = period_text
    set_font(run2, size=14, color=RGBColor(230, 230, 230))

    p3 = tf2.add_paragraph()
    run3 = p3.add_run()
    pm_name = data.get('pmName', '')
    run3.text = f"Project Manager: {pm_name}" if pm_name else report_date
    set_font(run3, size=12, color=RGBColor(200, 200, 200))

    logo_path = data.get('logoPath')
    if logo_path and os.path.exists(logo_path):
        try:
            slide.shapes.add_picture(logo_path, Inches(10.5), Inches(4.0), height=Inches(0.8))
        except Exception:
            pass

    return slide

def create_progress_summary_slide(prs, data, sections, primary_color, secondary_color):
    """Slide 2: Progress Summary with AI-generated narrative."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide, primary_color, top=0)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(10), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Progress Summary"
    set_font(run, size=24, bold=True, color=primary_color)

    metrics = data.get('metrics', {})
    if metrics:
        metrics_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.5))
        tf_m = metrics_box.text_frame
        p = tf_m.paragraphs[0]
        items = []
        if metrics.get('totalHours', '0') != '0':
            items.append(f"Hours: {metrics['totalHours']} ({metrics.get('billableHours', '0')} billable)")
        if metrics.get('teamMembers', 0) > 0:
            items.append(f"Team: {metrics['teamMembers']} members")
        if metrics.get('totalExpenses', '0.00') != '0.00':
            items.append(f"Expenses: ${metrics['totalExpenses']}")
        if items:
            run = p.add_run()
            run.text = "  |  ".join(items)
            set_font(run, size=9, color='#666666')

    content_top = Inches(1.6) if metrics else Inches(1.2)
    content_box = slide.shapes.add_textbox(Inches(0.8), content_top, Inches(11.5), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    summary_text = sections.get('Progress Summary', '')
    if summary_text:
        render_markdown_text(tf, summary_text, primary_color, size=11, start_fresh=True)

    # Milestone Posture intentionally omitted here — those details live on the
    # Timeline & Milestones slide, so showing them again would duplicate content.

    return slide

def _parse_task_string(task_str):
    """Parse a raw task allocation string into structured components."""
    name = task_str
    person = ''
    epic = ''
    stage = ''

    paren_match = re.match(r'^(.+?)\s*\(([^)]+)\)\s*[-–—]\s*(.+)$', task_str)
    if paren_match:
        name = paren_match.group(1).strip()
        person = paren_match.group(2).strip()
        rest = paren_match.group(3).strip()
    else:
        dash_match = re.match(r'^(.+?)\s*[-–—]\s*(.+)$', task_str)
        if dash_match:
            name = dash_match.group(1).strip()
            rest = dash_match.group(2).strip()
        else:
            rest = ''

    if rest:
        date_match = re.match(r'^(.+?)\s*\[.*$', rest)
        context = date_match.group(1).strip() if date_match else rest
        parts = context.split('>')
        if len(parts) >= 2:
            epic = parts[0].strip()
            stage = parts[1].strip()
        elif context:
            epic = context.strip()

    if not person:
        inner_match = re.match(r'^(.+?)\s*\(([^)]+)\)(.*)$', task_str)
        if inner_match:
            name = inner_match.group(1).strip()
            person = inner_match.group(2).strip()

    return {'name': name, 'person': person, 'epic': epic, 'stage': stage}


def _group_tasks_by_epic(prior_tasks, current_tasks):
    """Group task strings by epic for a cleaner fallback display."""
    groups = {}

    for task_str in prior_tasks:
        parsed = _parse_task_string(task_str)
        key = parsed['epic'] or 'General'
        if key not in groups:
            groups[key] = {'epic': key, 'tasks': []}
        groups[key]['tasks'].append({'name': parsed['name'], 'person': parsed['person'], 'done': True})

    for task_str in current_tasks:
        parsed = _parse_task_string(task_str)
        key = parsed['epic'] or 'General'
        if key not in groups:
            groups[key] = {'epic': key, 'tasks': []}
        groups[key]['tasks'].append({'name': parsed['name'], 'person': parsed['person'], 'done': False})

    result = sorted(groups.values(), key=lambda g: (-sum(1 for t in g['tasks'] if t['done']), g['epic']))
    return result


def create_accomplishments_slide(prs, data, sections, primary_color, secondary_color):
    """Slide 3: Key Accomplishments with rich AI narrative."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide, primary_color, top=0)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(10), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Key Accomplishments"
    set_font(run, size=24, bold=True, color=primary_color)

    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.5), Inches(6.0))
    tf = content_box.text_frame
    tf.word_wrap = True

    accomplishments_text = sections.get('Key Accomplishments', '')
    if accomplishments_text:
        items = parse_bullet_items(accomplishments_text)
        first = True
        for item in items:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()

            p.space_before = Pt(6)
            p.space_after = Pt(2)

            run = p.add_run()
            run.text = f"• {item['title']}"
            set_font(run, size=11, bold=True, color=primary_color)

            if item.get('description'):
                p2 = tf.add_paragraph()
                p2.space_before = Pt(1)
                p2.space_after = Pt(4)
                run2 = p2.add_run()
                run2.text = f"  {item['description']}"
                set_font(run2, size=10)

            for sub in item.get('sub_items', []):
                p3 = tf.add_paragraph()
                p3.space_before = Pt(1)
                run3 = p3.add_run()
                run3.text = f"    – {sub}"
                set_font(run3, size=9, color='#444444')
    else:
        print("[PPTX] No 'Key Accomplishments' section found in AI output — using grouped task fallback", file=sys.stderr)
        activities = data.get('projectActivities', {})
        prior = activities.get('prior', [])
        current = activities.get('current', [])
        has_fallback = prior or current
        if has_fallback:
            grouped = _group_tasks_by_epic(prior, current)
            first = True
            for group in grouped[:12]:
                if first:
                    p = tf.paragraphs[0]
                    first = False
                else:
                    p = tf.add_paragraph()
                p.space_before = Pt(6)
                p.space_after = Pt(2)
                run = p.add_run()
                run.text = f"• {group['epic']}"
                set_font(run, size=11, bold=True, color=primary_color)

                for task in group['tasks'][:6]:
                    p2 = tf.add_paragraph()
                    p2.space_before = Pt(1)
                    p2.space_after = Pt(1)
                    run2 = p2.add_run()
                    icon = '\u2713' if task['done'] else '\u25B8'
                    run2.text = f"  {icon} {task['name']}"
                    if task.get('person'):
                        run2.text += f" ({task['person']})"
                    set_font(run2, size=9, color='#333333' if task['done'] else '#555555')
                if len(group['tasks']) > 6:
                    p3 = tf.add_paragraph()
                    run3 = p3.add_run()
                    run3.text = f"    + {len(group['tasks']) - 6} more"
                    set_font(run3, size=8, italic=True, color='#888888')
        else:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "No accomplishments data available for this period."
            set_font(run, size=11, color='#888888')

    return slide

PRIORITY_COLORS = {
    'critical': '#DC2626',
    'high': '#EA580C',
    'medium': '#CA8A04',
    'low': '#16A34A',
}

STATUS_COLORS = {
    'open': '#3B82F6',
    'in_progress': '#F59E0B',
    'mitigated': '#8B5CF6',
    'resolved': '#22C55E',
    'closed': '#6B7280',
    'deferred': '#9CA3AF',
    'superseded': '#9CA3AF',
    'accepted': '#22C55E',
}

RAIDD_OPEN_STATUSES = ('open', 'in_progress')


def _classify_raidd(raidd_raw):
    """Classify flat raidd array into typed buckets."""
    buckets = {
        'risk': [],
        'action_item': [],
        'issue': [],
        'decision': [],
        'dependency': [],
    }
    if isinstance(raidd_raw, list):
        for entry in raidd_raw:
            t = entry.get('type', '')
            if t in buckets:
                buckets[t].append(entry)
    elif isinstance(raidd_raw, dict):
        buckets['risk'] = raidd_raw.get('risks', [])
        buckets['issue'] = raidd_raw.get('issues', [])
        buckets['action_item'] = raidd_raw.get('actionItems', [])
        buckets['decision'] = raidd_raw.get('decisions', [])
        buckets['dependency'] = raidd_raw.get('dependencies', [])
    return buckets


def _priority_sort_key(entry):
    order = {'critical': 0, 'high': 1, 'medium': 2, 'low': 3}
    return order.get(entry.get('priority', 'medium'), 2)


_RAIDD_SLIDE_HEIGHT_IN = 7.5
_RAIDD_ROW_HEIGHT_IN = 0.55
_RAIDD_BOTTOM_MARGIN_IN = 0.15
_RAIDD_TABLE_TOP_WITH_SUBTITLE_IN = 1.05
_RAIDD_TABLE_TOP_WITHOUT_SUBTITLE_IN = 0.85


def _raidd_max_rows(table_top_in):
    """Return the maximum number of data rows that fit on a slide given the table top offset."""
    available_in = _RAIDD_SLIDE_HEIGHT_IN - table_top_in - _RAIDD_BOTTOM_MARGIN_IN
    max_total_rows = int(available_in / _RAIDD_ROW_HEIGHT_IN)  # header + data rows
    return max(1, max_total_rows - 1)


def _add_raidd_table_slide(prs, title_text, subtitle_text, entries, columns, col_widths, primary_color, secondary_color, page_num=None, total_pages=None):
    """Create a single RAIDD category slide with a structured table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide, primary_color, top=0)

    page_label = f" ({page_num}/{total_pages})" if total_pages and total_pages > 1 else ""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(10), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{title_text}{page_label}"
    set_font(run, size=20, bold=True, color=primary_color)

    if subtitle_text and page_num in (1, None):
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.72), Inches(12), Inches(0.3))
        stf = sub_box.text_frame
        sp = stf.paragraphs[0]
        srun = sp.add_run()
        srun.text = subtitle_text
        set_font(srun, size=9, color='#888888')

    has_subtitle_header = bool(subtitle_text and page_num in (1, None))
    table_top_in = _RAIDD_TABLE_TOP_WITH_SUBTITLE_IN if has_subtitle_header else _RAIDD_TABLE_TOP_WITHOUT_SUBTITLE_IN
    table_top = Inches(table_top_in)

    # Compute max data rows dynamically so the table never overflows the slide boundary
    max_rows_per_page = _raidd_max_rows(table_top_in)
    display_entries = entries[:max_rows_per_page]
    rows = len(display_entries) + 1  # +1 for header
    cols = len(columns)

    # Cap table height to available slide space
    available_in = _RAIDD_SLIDE_HEIGHT_IN - table_top_in - _RAIDD_BOTTOM_MARGIN_IN
    table_height_in = min(_RAIDD_ROW_HEIGHT_IN * rows, available_in)
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.2), table_top, sum(col_widths), Inches(table_height_in))
    table = table_shape.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for i, (header, _) in enumerate(columns):
        set_cell_text(table.cell(0, i), header, size=10, bold=True, color='#FFFFFF', bg_color=primary_color, alignment=PP_ALIGN.LEFT)

    for row_idx, entry in enumerate(display_entries):
        r = row_idx + 1
        bg = '#FFFFFF' if r % 2 == 1 else '#F5F5FA'
        for col_idx, (_, extractor) in enumerate(columns):
            val, cell_opts = extractor(entry)
            color = cell_opts.get('color', None)
            bold = cell_opts.get('bold', False)
            set_cell_text(table.cell(r, col_idx), val, size=9.5, bold=bold, color=color, bg_color=bg, alignment=PP_ALIGN.LEFT)

    return slide




def create_raidd_slides(prs, data, sections, primary_color, secondary_color):
    """Generate RAIDD slides: a summary slide followed by per-category table slides."""
    raidd_raw = data.get('raidd', [])
    buckets = _classify_raidd(raidd_raw)

    open_statuses = RAIDD_OPEN_STATUSES
    decision_open_statuses = {'proposed', 'open', 'in_progress'}
    decision_closed_statuses = {'approved', 'rejected', 'closed', 'completed'}

    risks = sorted(buckets['risk'], key=_priority_sort_key)
    actions = sorted(buckets['action_item'], key=_priority_sort_key)
    issues = sorted(buckets['issue'], key=_priority_sort_key)
    all_decisions = buckets['decision']
    dependencies = sorted(buckets['dependency'], key=_priority_sort_key)

    _decision_filter_raw = data.get('decisionLogFilter', 'open')
    decision_log_filter = str(_decision_filter_raw).strip().lower() if _decision_filter_raw else 'open'
    if decision_log_filter == 'open':
        decisions = [d for d in all_decisions if d.get('status') in decision_open_statuses]
    elif decision_log_filter == 'closed':
        decisions = [d for d in all_decisions if d.get('status') in decision_closed_statuses]
    else:
        decisions = all_decisions

    active_risks = [e for e in risks if e.get('status') in open_statuses]
    active_issues = [e for e in issues if e.get('status') in open_statuses]
    active_actions = [e for e in actions if e.get('status') in open_statuses]
    active_deps = [e for e in dependencies if e.get('status') in open_statuses]

    # Counts and the critical summary follow the report's "active items only" runtime
    # setting: when limiting to active, show active items; otherwise show all items.
    # Parse defensively in case the flag ever arrives as a string ("false") rather than a bool.
    _open_only_raw = data.get('raiddOpenOnly', True)
    raidd_open_only = (
        _open_only_raw if isinstance(_open_only_raw, bool)
        else str(_open_only_raw).strip().lower() not in ('false', '0', 'no', 'off', '')
    )
    display_risks = active_risks if raidd_open_only else risks
    display_issues = active_issues if raidd_open_only else issues
    display_actions = active_actions if raidd_open_only else actions
    display_deps = active_deps if raidd_open_only else dependencies
    display_all = display_risks + display_issues + display_actions + display_deps

    critical_count = sum(1 for e in display_all if e.get('priority') == 'critical')
    high_count = sum(1 for e in display_all if e.get('priority') == 'high')

    overdue_actions = [e for e in active_actions if e.get('dueDate') and e['dueDate'] < datetime.now().strftime('%Y-%m-%d')]

    # --- RAIDD Summary Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide, primary_color, top=0)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(10), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "RAIDD Log Overview"
    set_font(run, size=22, bold=True, color=primary_color)

    def _secondary(active_n, total_n):
        if raidd_open_only or active_n == total_n:
            return None
        return f"{active_n} active"

    categories_summary = [
        ('Risks', len(display_risks), _secondary(len(active_risks), len(risks)), '#DC2626'),
        ('Action Items', len(display_actions), _secondary(len(active_actions), len(actions)), '#EA580C'),
        ('Issues', len(display_issues), _secondary(len(active_issues), len(issues)), '#CA8A04'),
        ('Decisions', len(decisions), f"{len(all_decisions) - len(decisions)} hidden" if decision_log_filter != 'all' and len(decisions) != len(all_decisions) else None, '#3B82F6'),
        ('Dependencies', len(display_deps), _secondary(len(active_deps), len(dependencies)), '#8B5CF6'),
    ]

    card_y = Inches(0.95)
    card_w = Inches(2.3)
    card_h = Inches(1.1)
    card_gap = Inches(0.25)
    start_x = Inches(0.5)

    for idx, (cat_name, active_count, secondary, accent) in enumerate(categories_summary):
        x = start_x + idx * (card_w + card_gap)
        card_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, card_y, card_w, card_h)
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = hex_to_rgb('#F8F8FC')
        card_bg.line.color.rgb = hex_to_rgb('#E5E5EA')
        card_bg.line.width = Pt(0.5)

        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, card_y, card_w, Inches(0.05))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = hex_to_rgb(accent)
        accent_bar.line.fill.background()

        num_box = slide.shapes.add_textbox(x + Inches(0.15), card_y + Inches(0.15), card_w - Inches(0.3), Inches(0.5))
        ntf = num_box.text_frame
        np = ntf.paragraphs[0]
        np.alignment = PP_ALIGN.LEFT
        nrun = np.add_run()
        nrun.text = str(active_count)
        set_font(nrun, size=28, bold=True, color=accent)

        label_box = slide.shapes.add_textbox(x + Inches(0.15), card_y + Inches(0.65), card_w - Inches(0.3), Inches(0.35))
        ltf = label_box.text_frame
        lp = ltf.paragraphs[0]
        lp.alignment = PP_ALIGN.LEFT
        lrun = lp.add_run()
        lrun.text = f"{cat_name}"
        set_font(lrun, size=10, bold=True, color='#333333')
        if secondary:
            lrun2 = lp.add_run()
            lrun2.text = f"  ({secondary})"
            set_font(lrun2, size=8, color='#888888')

    alert_y = card_y + card_h + Inches(0.3)
    alerts = []
    if critical_count > 0:
        alerts.append(f"⚠ {critical_count} CRITICAL item{'s' if critical_count != 1 else ''} require immediate attention")
    if high_count > 0:
        alerts.append(f"▲ {high_count} HIGH priority item{'s' if high_count != 1 else ''}")
    if overdue_actions:
        alerts.append(f"⏰ {len(overdue_actions)} overdue action item{'s' if len(overdue_actions) != 1 else ''}")

    if alerts:
        alert_box = slide.shapes.add_textbox(Inches(0.5), alert_y, Inches(12), Inches(0.6))
        atf = alert_box.text_frame
        atf.word_wrap = True
        for i, alert_text in enumerate(alerts):
            ap = atf.paragraphs[0] if i == 0 else atf.add_paragraph()
            ap.space_before = Pt(2)
            arun = ap.add_run()
            arun.text = alert_text
            color = '#DC2626' if 'CRITICAL' in alert_text else ('#EA580C' if 'HIGH' in alert_text else '#B45309')
            set_font(arun, size=10, bold=True, color=color)
        alert_y += Inches(0.15) * len(alerts) + Inches(0.3)

    # Summary page lists CRITICAL items only — the per-category detail slides that
    # follow contain every open item regardless of criticality.
    type_labels = {
        'risk': 'Risk', 'issue': 'Issue', 'action_item': 'Action Item', 'dependency': 'Dependency',
    }
    critical_groups = [
        ('risk', display_risks),
        ('issue', display_issues),
        ('action_item', display_actions),
        ('dependency', display_deps),
    ]
    critical_items = [
        (type_key, e)
        for type_key, group in critical_groups
        for e in group
        if e.get('priority') == 'critical'
    ]

    crit_y = alert_y if alerts else card_y + card_h + Inches(0.3)
    # Box extends to 0.15" from slide bottom so content is always within bounds.
    crit_box_h = max(Inches(0.5), Inches(7.35) - crit_y)
    crit_box = slide.shapes.add_textbox(Inches(0.5), crit_y, Inches(12.3), crit_box_h)
    ctf = crit_box.text_frame
    ctf.word_wrap = True

    header_p = ctf.paragraphs[0]
    header_run = header_p.add_run()
    header_run.text = "Critical Items Requiring Attention"
    set_font(header_run, size=11, bold=True, color=primary_color)

    # Cap items so they always fit at 8pt minimum: each item occupies ~0.35"
    # (one title line + one sub line + spacing). The available height minus the
    # 0.2" header leaves room for at most this many items.
    _crit_avail_in = (crit_box_h / 914400) - 0.2
    MAX_CRIT_ITEMS = max(4, int(_crit_avail_in / 0.35))

    if critical_items:
        shown = critical_items[:MAX_CRIT_ITEMS]
        hidden = len(critical_items) - len(shown)
        for type_key, e in shown:
            ref = (e.get('refNumber', '') or '').strip()
            title = (e.get('title', '') or '').strip()
            label = type_labels.get(type_key, type_key.replace('_', ' ').title())

            ip = ctf.add_paragraph()
            ip.space_before = Pt(4)
            ref_run = ip.add_run()
            ref_run.text = f"{ref + ' ' if ref else ''}{title}  "
            set_font(ref_run, size=9, bold=True, color='#222222')
            tag_run = ip.add_run()
            tag_run.text = f"[{label}]"
            set_font(tag_run, size=8, bold=True, color=PRIORITY_COLORS.get('critical', '#DC2626'))

            detail = (e.get('mitigationPlan', '') or '').strip()
            meta_parts = []
            owner = (e.get('ownerName', '') or '').strip()
            due = (e.get('dueDate', '') or '').strip()
            if owner:
                meta_parts.append(f"Owner: {owner}")
            if due:
                meta_parts.append(f"Due: {due}")
            sub_text = detail
            if meta_parts:
                sub_text = (sub_text + '  ' if sub_text else '') + ' | '.join(meta_parts)
            if sub_text:
                sp = ctf.add_paragraph()
                sp.space_before = Pt(1)
                sub_run = sp.add_run()
                sub_run.text = sub_text
                set_font(sub_run, size=8, color='#555555')

        if hidden > 0:
            mp = ctf.add_paragraph()
            mp.space_before = Pt(6)
            mrun = mp.add_run()
            mrun.text = f"… and {hidden} more critical item{'s' if hidden != 1 else ''} — see detail slides"
            set_font(mrun, size=8, italic=True, color='#888888')
    else:
        np = ctf.add_paragraph()
        np.space_before = Pt(4)
        nrun = np.add_run()
        nrun.text = "No critical items at this time. See the following slides for all open RAIDD items."
        set_font(nrun, size=9, color='#555555')

    # --- Per-category detail slides ---
    # Use worst-case table_top (page 1 with subtitle) so pagination is consistent
    # across all pages — later pages without the subtitle header have slightly more
    # room but using the tighter bound keeps page counts stable and avoids overflow.
    max_rows_per_page = _raidd_max_rows(_RAIDD_TABLE_TOP_WITH_SUBTITLE_IN)

    def _priority_cell(entry):
        p = entry.get('priority', '')
        return (p.upper() if p else '—', {'color': PRIORITY_COLORS.get(p, '#888888'), 'bold': True})

    def _status_cell(entry):
        s = entry.get('status', '')
        label = s.replace('_', ' ').title() if s else '—'
        return (label, {'color': STATUS_COLORS.get(s, '#888888'), 'bold': True})

    def _owner_cell(entry):
        return (entry.get('ownerName', '') or entry.get('assigneeName', '') or '—', {})

    def _ref_cell(entry):
        return (entry.get('refNumber', '') or '—', {'bold': True})

    def _title_cell(entry):
        return ((entry.get('title', '') or '').strip(), {})

    def _due_cell(entry):
        d = entry.get('dueDate', '')
        if d and d < datetime.now().strftime('%Y-%m-%d') and entry.get('status') in RAIDD_OPEN_STATUSES:
            return (d, {'color': '#DC2626', 'bold': True})
        return (d or '—', {})

    def _desc_cell(entry):
        return ((entry.get('description', '') or entry.get('mitigationPlan', '') or '').strip(), {})

    def _impact_cell(entry):
        imp = entry.get('impact', '')
        return (imp.title() if imp else '—', {'color': PRIORITY_COLORS.get(imp, '#888888')})

    def _likelihood_cell(entry):
        lk = entry.get('likelihood', '')
        return (lk.replace('_', ' ').title() if lk else '—', {})

    def _mitigation_cell(entry):
        return ((entry.get('mitigationPlan', '') or '').strip(), {})

    def _resolution_cell(entry):
        text = (entry.get('resolutionNotes', '') or entry.get('description', '') or '').strip()
        return (text, {})

    def _paginate_category(entries, title, subtitle, columns, col_widths):
        if not entries:
            return
        pages = []
        for i in range(0, len(entries), max_rows_per_page):
            pages.append(entries[i:i + max_rows_per_page])

        total_pages = len(pages)
        for page_idx, page_entries in enumerate(pages):
            _add_raidd_table_slide(
                prs, title, subtitle, page_entries,
                columns, col_widths, primary_color, secondary_color,
                page_num=page_idx + 1, total_pages=total_pages,
            )

    # Risks table
    if risks:
        active_label = f"{len(active_risks)} active" if active_risks else "None active"
        _paginate_category(
            risks,
            "Risks",
            f"{len(risks)} total  |  {active_label}  |  {sum(1 for r in active_risks if r.get('priority') == 'critical')} critical",
            [
                ('Ref', _ref_cell),
                ('Risk', _title_cell),
                ('Priority', _priority_cell),
                ('Impact', _impact_cell),
                ('Likelihood', _likelihood_cell),
                ('Status', _status_cell),
                ('Owner', _owner_cell),
                ('Mitigation', _mitigation_cell),
            ],
            [Inches(0.7), Inches(2.8), Inches(0.9), Inches(0.8), Inches(1.0), Inches(1.0), Inches(1.3), Inches(4.2)],
        )

    # Action Items table
    if actions:
        overdue_label = f"  |  {len(overdue_actions)} overdue" if overdue_actions else ""
        _paginate_category(
            actions,
            "Action Items",
            f"{len(actions)} total  |  {len(active_actions)} open{overdue_label}",
            [
                ('Ref', _ref_cell),
                ('Action', _title_cell),
                ('Priority', _priority_cell),
                ('Status', _status_cell),
                ('Assignee', _owner_cell),
                ('Due Date', _due_cell),
                ('Notes', _desc_cell),
            ],
            [Inches(0.7), Inches(3.0), Inches(0.9), Inches(1.0), Inches(1.5), Inches(1.1), Inches(4.5)],
        )

    # Issues table
    if issues:
        _paginate_category(
            issues,
            "Issues",
            f"{len(issues)} total  |  {len(active_issues)} active",
            [
                ('Ref', _ref_cell),
                ('Issue', _title_cell),
                ('Priority', _priority_cell),
                ('Status', _status_cell),
                ('Owner', _owner_cell),
                ('Resolution', _resolution_cell),
            ],
            [Inches(0.7), Inches(3.2), Inches(0.9), Inches(1.0), Inches(1.5), Inches(5.4)],
        )

    # Decisions table
    if decisions:
        if decision_log_filter == 'open':
            decisions_subtitle = f"{len(decisions)} open  |  {len(all_decisions)} total"
        elif decision_log_filter == 'closed':
            decisions_subtitle = f"{len(decisions)} closed  |  {len(all_decisions)} total"
        else:
            open_count = sum(1 for d in all_decisions if d.get('status') in decision_open_statuses)
            decisions_subtitle = f"{len(decisions)} total  |  {open_count} open"
        _paginate_category(
            decisions,
            "Decisions",
            decisions_subtitle,
            [
                ('Ref', _ref_cell),
                ('Decision', _title_cell),
                ('Status', _status_cell),
                ('Owner', _owner_cell),
                ('Resolution / Rationale', _resolution_cell),
            ],
            [Inches(0.7), Inches(3.5), Inches(1.0), Inches(1.5), Inches(6.0)],
        )

    # Dependencies table
    if dependencies:
        _paginate_category(
            dependencies,
            "Dependencies",
            f"{len(dependencies)} total  |  {len(active_deps)} active",
            [
                ('Ref', _ref_cell),
                ('Dependency', _title_cell),
                ('Priority', _priority_cell),
                ('Status', _status_cell),
                ('Owner', _owner_cell),
                ('Due Date', _due_cell),
                ('Notes', _desc_cell),
            ],
            [Inches(0.7), Inches(3.0), Inches(0.9), Inches(1.0), Inches(1.5), Inches(1.1), Inches(4.5)],
        )

def create_upcoming_slide(prs, data, sections, primary_color, secondary_color):
    """Slide 5: Upcoming Activities with AI narrative."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide, primary_color, top=0)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(10), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Upcoming Activities"
    set_font(run, size=24, bold=True, color=primary_color)

    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.5), Inches(6.0))
    tf = content_box.text_frame
    tf.word_wrap = True

    upcoming_text = sections.get('Upcoming Activities', '')
    if upcoming_text:
        items = parse_bullet_items(upcoming_text)
        first = True
        for item in items:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()

            p.space_before = Pt(6)
            p.space_after = Pt(2)

            if item['title']:
                run = p.add_run()
                run.text = f"• {item['title']}"
                set_font(run, size=11, bold=True, color=primary_color)
            else:
                # No bold-prefixed title — render full line with inline bold support
                render_inline_bold(p, f"• {item.get('description', '')}", size=11, primary_color=primary_color)

            if item['title'] and item.get('description'):
                p2 = tf.add_paragraph()
                p2.space_before = Pt(1)
                p2.space_after = Pt(4)
                render_inline_bold(p2, f"  {item['description']}", size=10)

            for sub in item.get('sub_items', []):
                p3 = tf.add_paragraph()
                p3.space_before = Pt(1)
                run3 = p3.add_run()
                run3.text = f"    – {sub}"
                set_font(run3, size=9, color='#444444')
    else:
        activities = data.get('projectActivities', {})
        upcoming = activities.get('upcoming', [])
        if upcoming:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"Scheduled Tasks ({len(upcoming)})"
            set_font(run, size=13, bold=True, color=primary_color)
            for act in upcoming[:18]:
                p2 = tf.add_paragraph()
                p2.space_before = Pt(3)
                p2.space_after = Pt(2)
                run2 = p2.add_run()
                run2.text = f"  \u25B8 {act}"
                set_font(run2, size=9)
            if len(upcoming) > 18:
                p2 = tf.add_paragraph()
                run2 = p2.add_run()
                run2.text = f"    ... and {len(upcoming) - 18} more scheduled tasks"
                set_font(run2, size=9, italic=True, color='#666666')
        else:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "No upcoming activities data available."
            set_font(run, size=11, color='#888888')

    return slide

def create_timeline_slide(prs, data, primary_color, secondary_color):
    """Slide 6: Timeline & Milestones - Gantt chart with native PowerPoint shapes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide, primary_color, top=0)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(10), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Timeline & Milestones"
    set_font(run, size=24, bold=True, color=primary_color)

    timeline = data.get('timeline', {})
    epic_groups = timeline.get('epicGroups', [])
    unlinked_milestones = timeline.get('unlinkedMilestones', [])
    payment_milestones = timeline.get('paymentMilestones', [])
    all_milestones = data.get('milestones', [])

    has_epic_data = any(eg.get('stages') for eg in epic_groups)
    has_any_data = has_epic_data or unlinked_milestones or payment_milestones or all_milestones

    if not has_any_data:
        note_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(1))
        tf = note_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "No timeline data available for this project."
        set_font(run, size=12, color='#666666')
        return slide

    all_dates = []
    for eg in epic_groups:
        for stage in eg.get('stages', []):
            if stage.get('startDate'):
                all_dates.append(stage['startDate'])
            if stage.get('endDate'):
                all_dates.append(stage['endDate'])
        for ms in eg.get('milestones', []):
            if ms.get('targetDate'):
                all_dates.append(ms['targetDate'])
    for ms in unlinked_milestones:
        if ms.get('targetDate'):
            all_dates.append(ms['targetDate'])
    for ms in payment_milestones:
        if ms.get('targetDate'):
            all_dates.append(ms['targetDate'])

    has_gantt_dates = len(all_dates) > 0

    if has_gantt_dates:
        all_dates.sort()
        min_date = datetime.strptime(all_dates[0][:10], '%Y-%m-%d')
        max_date = datetime.strptime(all_dates[-1][:10], '%Y-%m-%d')

        padding_days = max(7, int((max_date - min_date).days * 0.05))
        chart_start = min_date - timedelta(days=padding_days)
        chart_end = max_date + timedelta(days=padding_days)
        total_days = max((chart_end - chart_start).days, 1)

        label_width = Inches(2.8)
        chart_left = Inches(3.2)
        chart_width_in = 9.5
        chart_width = Inches(chart_width_in)
        chart_top = Inches(1.3)
        row_height = Inches(0.35)
        epic_header_height = Inches(0.30)
        milestone_dot_size = Inches(0.18)

        def date_to_x(date_str):
            d = datetime.strptime(date_str[:10], '%Y-%m-%d')
            frac = (d - chart_start).days / total_days
            return chart_left + Emu(int(chart_width_in * frac * 914400))

        bar_colors = ['#6366f1', '#8b5cf6', '#a78bfa', '#7c3aed', '#4f46e5', '#818cf8', '#c084fc', '#a855f7']

        _draw_month_axis(slide, chart_start, chart_end, chart_left, chart_width, chart_top, total_days, primary_color)

        y_cursor = chart_top + Inches(0.35)
        color_idx = 0

        for eg in epic_groups:
            epic_name = eg.get('epicName', 'Unnamed Epic')
            stages = eg.get('stages', [])
            ms_list = eg.get('milestones', [])

            if not stages and not ms_list:
                continue

            epic_label = slide.shapes.add_textbox(Inches(0.4), y_cursor, label_width, epic_header_height)
            tf = epic_label.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = epic_name
            set_font(run, size=9, bold=True, color=primary_color)
            y_cursor += epic_header_height

            for stage in stages:
                start_str = stage.get('startDate', '')
                end_str = stage.get('endDate', '')

                stage_label = slide.shapes.add_textbox(Inches(0.6), y_cursor, label_width - Inches(0.2), row_height)
                tf = stage_label.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.RIGHT
                run = p.add_run()
                stage_name = stage.get('name', '')
                if start_str and end_str:
                    run.text = stage_name
                else:
                    run.text = f"{stage_name}  (no dates)"
                set_font(run, size=8, color='#333333')

                if start_str and end_str:
                    bar_left = date_to_x(start_str)
                    bar_right = date_to_x(end_str)
                    bar_w = max(bar_right - bar_left, Emu(int(914400 * 0.1)))
                    bar_h = Inches(0.22)
                    bar_top = y_cursor + Inches(0.06)

                    bar_color = bar_colors[color_idx % len(bar_colors)]
                    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, bar_top, bar_w, bar_h)
                    bar.fill.solid()
                    bar.fill.fore_color.rgb = hex_to_rgb(bar_color)
                    bar.line.fill.background()

                    if bar_w > Emu(int(914400 * 1.2)):
                        bar_tf = bar.text_frame
                        bar_tf.word_wrap = False
                        bp = bar_tf.paragraphs[0]
                        bp.alignment = PP_ALIGN.CENTER
                        brun = bp.add_run()
                        s = datetime.strptime(start_str[:10], '%Y-%m-%d')
                        e = datetime.strptime(end_str[:10], '%Y-%m-%d')
                        brun.text = f"{s.strftime('%b %d')} – {e.strftime('%b %d')}"
                        set_font(brun, size=6, color='#FFFFFF', bold=True)

                y_cursor += row_height
                color_idx += 1

            for ms in ms_list:
                if not ms.get('targetDate'):
                    continue
                _draw_milestone_dot(slide, ms, date_to_x, y_cursor, label_width, secondary_color, milestone_dot_size)
                y_cursor += Inches(0.28)

        if unlinked_milestones:
            ul_label = slide.shapes.add_textbox(Inches(0.4), y_cursor, label_width, epic_header_height)
            tf = ul_label.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Project Milestones"
            set_font(run, size=9, bold=True, color=primary_color)
            y_cursor += epic_header_height

            for ms in unlinked_milestones:
                if not ms.get('targetDate'):
                    continue
                _draw_milestone_dot(slide, ms, date_to_x, y_cursor, label_width, secondary_color, milestone_dot_size)
                y_cursor += Inches(0.28)

        if payment_milestones:
            _draw_payment_milestones_section(slide, payment_milestones, y_cursor, primary_color, secondary_color, date_to_x if has_gantt_dates else None)

    else:
        y_cursor = Inches(1.3)
        _draw_epic_stages_list(slide, epic_groups, y_cursor, primary_color)
        y_cursor_after = Inches(1.3) + Inches(0.3) * sum(1 + len(eg.get('stages', [])) for eg in epic_groups if eg.get('stages'))
        if payment_milestones:
            _draw_payment_milestones_section(slide, payment_milestones, y_cursor_after, primary_color, secondary_color, None)

    return slide


def _draw_month_axis(slide, chart_start, chart_end, chart_left, chart_width, chart_top, total_days, primary_color):
    """Draw month labels along the top axis of the Gantt chart."""
    chart_width_in = 9.5

    axis_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, chart_left, chart_top + Inches(0.28), chart_width, Inches(0.01))
    axis_line.fill.solid()
    axis_line.fill.fore_color.rgb = hex_to_rgb('#cccccc')
    axis_line.line.fill.background()

    current = datetime(chart_start.year, chart_start.month, 1)
    while current <= chart_end:
        if current >= chart_start:
            frac = (current - chart_start).days / total_days
            x = chart_left + Emu(int(chart_width_in * frac * 914400))

            tick = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, chart_top + Inches(0.25), Inches(0.01), Inches(0.06))
            tick.fill.solid()
            tick.fill.fore_color.rgb = hex_to_rgb('#aaaaaa')
            tick.line.fill.background()

            label_box = slide.shapes.add_textbox(x - Inches(0.3), chart_top, Inches(0.8), Inches(0.25))
            tf = label_box.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = current.strftime('%b %Y')
            set_font(run, size=7, color='#666666')

        if current.month == 12:
            current = datetime(current.year + 1, 1, 1)
        else:
            current = datetime(current.year, current.month + 1, 1)


def _draw_milestone_dot(slide, ms, date_to_x, y_cursor, label_width, secondary_color, dot_size):
    """Draw a single milestone as a diamond dot with label."""
    target = ms.get('targetDate', '')
    if not target:
        return

    x = date_to_x(target)
    dot_top = y_cursor + Inches(0.04)

    status = ms.get('status', '')
    if status == 'completed':
        dot_color = '#22c55e'
    elif status == 'in-progress':
        dot_color = '#3b82f6'
    else:
        dot_color = secondary_color

    dot = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, x - dot_size // 2, dot_top, dot_size, dot_size)
    dot.fill.solid()
    dot.fill.fore_color.rgb = hex_to_rgb(dot_color)
    dot.line.fill.background()

    ms_label = slide.shapes.add_textbox(Inches(0.6), y_cursor, label_width - Inches(0.2), Inches(0.25))
    tf = ms_label.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    ms_name = ms.get('name', '')
    is_payment = ms.get('isPayment', False)
    label_suffix = ' $' if is_payment else ''
    d = datetime.strptime(target[:10], '%Y-%m-%d')
    run.text = f"{ms_name}{label_suffix} ({d.strftime('%b %d')})"
    set_font(run, size=7, color='#555555', italic=True)


def _draw_payment_milestones_section(slide, payment_milestones, y_cursor, primary_color, secondary_color, date_to_x_fn):
    """Draw payment milestones section below the Gantt chart."""
    if not payment_milestones:
        return y_cursor

    label_width = Inches(2.8)
    epic_header_height = Inches(0.30)
    row_height = Inches(0.28)

    header = slide.shapes.add_textbox(Inches(0.4), y_cursor, Inches(4), epic_header_height)
    tf = header.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Payment Milestones"
    set_font(run, size=9, bold=True, color=primary_color)
    y_cursor += epic_header_height

    for ms in payment_milestones:
        ms_name = ms.get('name', '')
        status = ms.get('status', '')
        target_date = ms.get('targetDate', '')

        status_icon = ''
        dot_color = secondary_color
        if status == 'completed':
            status_icon = ' \u2713'
            dot_color = '#22c55e'
        elif status == 'in-progress':
            status_icon = ' \u25B6'
            dot_color = '#3b82f6'

        ms_label = slide.shapes.add_textbox(Inches(0.6), y_cursor, label_width - Inches(0.2), Inches(0.25))
        tf = ms_label.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        date_display = ''
        if target_date:
            try:
                d = datetime.strptime(target_date[:10], '%Y-%m-%d')
                date_display = f" ({d.strftime('%b %d')})"
            except:
                pass
        run.text = f"$ {ms_name}{status_icon}{date_display}"
        set_font(run, size=8, color='#333333', bold=True)

        if target_date and date_to_x_fn:
            dot_size = Inches(0.16)
            x = date_to_x_fn(target_date)
            dot_top = y_cursor + Inches(0.04)
            dot = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, x - dot_size // 2, dot_top, dot_size, dot_size)
            dot.fill.solid()
            dot.fill.fore_color.rgb = hex_to_rgb(dot_color)
            dot.line.fill.background()
        else:
            status_txt = status.replace('-', ' ').title() if status else 'Pending'
            status_label = slide.shapes.add_textbox(Inches(3.2), y_cursor, Inches(2), Inches(0.25))
            tf = status_label.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = status_txt
            font_color = '#22c55e' if status == 'completed' else '#3b82f6' if status == 'in-progress' else '#9ca3af'
            set_font(run, size=7, color=font_color, italic=True)

        y_cursor += row_height

    return y_cursor


def _draw_epic_stages_list(slide, epic_groups, y_cursor, primary_color):
    """Draw a simple list of epics and their stages when no date data is available for Gantt bars."""
    label_width = Inches(10)
    epic_header_height = Inches(0.30)
    row_height = Inches(0.25)

    for eg in epic_groups:
        epic_name = eg.get('epicName', 'Unnamed Epic')
        stages = eg.get('stages', [])
        if not stages:
            continue

        epic_label = slide.shapes.add_textbox(Inches(0.6), y_cursor, label_width, epic_header_height)
        tf = epic_label.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = epic_name
        set_font(run, size=10, bold=True, color=primary_color)
        y_cursor += epic_header_height

        for stage in stages:
            stage_label = slide.shapes.add_textbox(Inches(1.0), y_cursor, label_width, row_height)
            tf = stage_label.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"\u2022  {stage.get('name', '')}"
            set_font(run, size=9, color='#444444')
            y_cursor += row_height

    return y_cursor


def _draw_milestone_table_fallback(slide, milestones, primary_color):
    """Fallback: draw a simple milestone table when no Gantt data is available."""
    if not milestones:
        note_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(1))
        tf = note_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "No milestones defined for this project."
        set_font(run, size=12, color='#666666')
        return

    num_rows = min(len(milestones) + 1, 16)
    table = slide.shapes.add_table(num_rows, 4, Inches(0.4), Inches(1.1), Inches(12.4), Inches(0.35 * num_rows)).table
    table.columns[0].width = Inches(5.0)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.5)
    table.columns[3].width = Inches(2.4)

    headers = ['Milestone', 'Target Date', 'Status', 'Dates']
    for i, header in enumerate(headers):
        set_cell_text(table.cell(0, i), header, size=9, bold=True, bg_color=primary_color, color=RGBColor(255, 255, 255))

    for row_idx, ms in enumerate(milestones[:15]):
        r = row_idx + 1
        set_cell_text(table.cell(r, 0), ms.get('name', ''), size=9)
        set_cell_text(table.cell(r, 1), ms.get('targetDate', ''), size=9)

        status = ms.get('status', '')
        status_color = None
        if status == 'completed':
            status_color = '#22c55e'
        elif status == 'in-progress':
            status_color = '#3b82f6'
        elif status == 'not-started':
            status_color = '#9ca3af'
        set_cell_text(table.cell(r, 2), status.replace('-', ' ').title(), size=9, color=status_color)

        dates = ''
        if ms.get('startDate') and ms.get('endDate'):
            dates = f"{ms['startDate']} – {ms['endDate']}"
        elif ms.get('startDate'):
            dates = ms['startDate']
        set_cell_text(table.cell(r, 3), dates, size=9)

def create_project_plan_slides(prs, data, primary_color, secondary_color):
    """Project Plan slides: assignments grouped by epic → stage, sorted by start date.
    Single line per assignment. Auto-paginates across multiple slides."""
    project_plan = data.get('projectPlan')
    if not project_plan:
        return

    groups = project_plan.get('groups', [])
    plan_filter = project_plan.get('filter', 'open')
    if not groups:
        return

    PAGE_TOP = Inches(1.1)
    PAGE_BOTTOM = Inches(7.0)
    LEFT_MARGIN = Inches(0.5)
    CONTENT_WIDTH = Inches(12.3)
    EPIC_HEIGHT = Inches(0.32)
    STAGE_HEIGHT = Inches(0.28)
    ROW_HEIGHT = Inches(0.35)
    HEADER_HEIGHT = Inches(0.28)

    STATUS_ICONS = {
        'open': '\u25CB',
        'in_progress': '\u25D2',
        'completed': '\u2713',
        'cancelled': '\u2717',
        'canceled': '\u2717',
        'superseded': '\u2717',
        'obsolete': '\u2717',
    }
    STATUS_COLORS = {
        'open': '#9ca3af',
        'in_progress': '#3b82f6',
        'completed': '#22c55e',
        'cancelled': '#ef4444',
        'canceled': '#ef4444',
        'superseded': '#ef4444',
        'obsolete': '#ef4444',
    }

    flat_items = []
    for group in groups:
        epic_name = group.get('epicName', 'Unnamed Epic')
        flat_items.append(('epic', epic_name, None))
        for stage_data in group.get('stages', []):
            stage_name = stage_data.get('stageName', '')
            flat_items.append(('stage', stage_name, None))
            for assignment in stage_data.get('assignments', []):
                flat_items.append(('assignment', None, assignment))

    def item_height(item_type):
        if item_type == 'epic':
            return EPIC_HEIGHT
        elif item_type == 'stage':
            return STAGE_HEIGHT
        else:
            return ROW_HEIGHT

    def start_new_slide(page_num):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_accent_bar(slide, primary_color, top=0)

        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(10), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        filter_label = "All Assignments" if plan_filter == 'all' else "Open Assignments"
        suffix = f" (Page {page_num})" if page_num > 1 else ""
        run.text = f"Project Plan — {filter_label}{suffix}"
        set_font(run, size=22, bold=True, color=primary_color)

        y = PAGE_TOP
        cols = [
            (Inches(0.3), 'Status'),
            (Inches(2.8), 'Assignee'),
            (Inches(3.8), 'Task'),
            (Inches(1.3), 'Hours'),
            (Inches(1.5), 'Start'),
            (Inches(1.5), 'End'),
        ]
        x = LEFT_MARGIN
        for col_w, col_label in cols:
            hdr = slide.shapes.add_textbox(x, y, col_w, HEADER_HEIGHT)
            tf_h = hdr.text_frame
            tf_h.word_wrap = False
            p_h = tf_h.paragraphs[0]
            run_h = p_h.add_run()
            run_h.text = col_label
            set_font(run_h, size=7, bold=True, color='#666666')
            x += col_w

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, LEFT_MARGIN, y + HEADER_HEIGHT - Inches(0.02), CONTENT_WIDTH, Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb('#dddddd')
        line.line.fill.background()

        return slide, y + HEADER_HEIGHT

    page_num = 1
    slide, y_cursor = start_new_slide(page_num)

    for item_type, label, assignment in flat_items:
        h = item_height(item_type)
        if y_cursor + h > PAGE_BOTTOM:
            page_num += 1
            slide, y_cursor = start_new_slide(page_num)

        if item_type == 'epic':
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, LEFT_MARGIN, y_cursor, CONTENT_WIDTH, EPIC_HEIGHT)
            bg.fill.solid()
            bg.fill.fore_color.rgb = hex_to_rgb(primary_color)
            bg.line.fill.background()

            epic_txt = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.15), y_cursor, Inches(10), EPIC_HEIGHT)
            tf = epic_txt.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = label
            set_font(run, size=9, bold=True, color=RGBColor(255, 255, 255))
            y_cursor += EPIC_HEIGHT

        elif item_type == 'stage':
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, LEFT_MARGIN, y_cursor, CONTENT_WIDTH, STAGE_HEIGHT)
            bg.fill.solid()
            bg.fill.fore_color.rgb = hex_to_rgb('#f0f0f0')
            bg.line.fill.background()

            stage_txt = slide.shapes.add_textbox(LEFT_MARGIN + Inches(0.3), y_cursor, Inches(10), STAGE_HEIGHT)
            tf = stage_txt.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = label
            set_font(run, size=8, bold=True, color='#333333')
            y_cursor += STAGE_HEIGHT

        else:
            status = assignment.get('status', 'open')
            icon = STATUS_ICONS.get(status, '\u25CB')
            icon_color = STATUS_COLORS.get(status, '#9ca3af')

            assignee = assignment.get('assignee', '')
            task = assignment.get('task', '')
            hours = assignment.get('hours', 0)
            start_d = assignment.get('startDate', '')
            end_d = assignment.get('endDate', '')

            if start_d:
                try:
                    sd = datetime.strptime(start_d[:10], '%Y-%m-%d')
                    start_d = sd.strftime('%b %d, %Y')
                except:
                    pass
            if end_d:
                try:
                    ed = datetime.strptime(end_d[:10], '%Y-%m-%d')
                    end_d = ed.strftime('%b %d, %Y')
                except:
                    pass

            x = LEFT_MARGIN
            cols_data = [
                (Inches(0.3), icon, icon_color, True),
                (Inches(2.8), assignee, '#222222', False),
                (Inches(3.8), task, '#444444', False),
                (Inches(1.3), f"{hours:.1f}h" if hours else '', '#555555', False),
                (Inches(1.5), start_d, '#555555', False),
                (Inches(1.5), end_d, '#555555', False),
            ]
            for col_w, text, color, is_icon in cols_data:
                cell = slide.shapes.add_textbox(x, y_cursor, col_w, ROW_HEIGHT)
                tf = cell.text_frame
                tf.word_wrap = not is_icon
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = str(text)
                sz = 10 if is_icon else 8
                set_font(run, size=sz, color=color)
                x += col_w

            y_cursor += ROW_HEIGHT


def create_deliverables_slide(prs, data, primary_color, secondary_color):
    """Deliverables tracking — max 10 items per slide, overflow to additional slides.
    All fonts are kept at 8pt minimum. Phase/stage headers carry over to new pages
    when a group spans the page boundary."""
    deliverables = data.get('deliverables', [])
    if not deliverables:
        return None

    ITEMS_PER_PAGE = 10  # max deliverable rows per slide (phase headers don't count)
    MIN_FONT = 8         # never go below this

    total = len(deliverables)
    accepted    = sum(1 for d in deliverables if d.get('status') == 'accepted')
    in_progress = sum(1 for d in deliverables if d.get('status') == 'in-progress')
    in_review   = sum(1 for d in deliverables if d.get('status') == 'in-review')
    not_started = sum(1 for d in deliverables if d.get('status') == 'not-started')
    rejected    = sum(1 for d in deliverables if d.get('status') == 'rejected')

    summary_parts = [f"{total} Total"]
    if accepted:    summary_parts.append(f"{accepted} Accepted")
    if in_review:   summary_parts.append(f"{in_review} In Review")
    if in_progress: summary_parts.append(f"{in_progress} In Progress")
    if not_started: summary_parts.append(f"{not_started} Not Started")
    if rejected:    summary_parts.append(f"{rejected} Rejected")
    summary_text = "  |  ".join(summary_parts)

    status_colors = {
        'accepted':    '#22C55E',
        'in-review':   '#3B82F6',
        'in-progress': '#F59E0B',
        'not-started': '#9CA3AF',
        'rejected':    '#EF4444',
    }

    cols = 5
    col_widths = [Inches(4.0), Inches(2.0), Inches(1.8), Inches(2.5), Inches(2.5)]

    # --- Group deliverables by stage, unassigned last ---
    phase_order = []
    phase_groups = {}
    for d in deliverables:
        key = d.get('stageId') or '__none__'
        if key not in phase_groups:
            stage_name = d.get('stageName') or 'Unassigned'
            epic_name  = d.get('epicName') or ''
            header = f"{epic_name} — {stage_name}" if (epic_name and key != '__none__') else stage_name
            phase_groups[key] = {'name': header, 'sort_name': stage_name, 'items': []}
            phase_order.append(key)
    for d in deliverables:
        phase_groups[d.get('stageId') or '__none__']['items'].append(d)

    # Sort groups by stage name using natural numeric sort so "1-Discovery",
    # "3-Prototype", "4-Pilot", "5-Handover" appear in numeric sequence
    # rather than insertion order from the DB.
    import re as _re
    def _phase_sort_key(k):
        name = phase_groups[k].get('sort_name', '')
        return [int(c) if c.isdigit() else c.lower() for c in _re.split(r'(\d+)', name)]

    named_keys = sorted([k for k in phase_order if k != '__none__'], key=_phase_sort_key)
    phase_order = named_keys + (['__none__'] if '__none__' in phase_groups else [])

    # --- Build full flat body (no cap) ---
    full_body = []
    for key in phase_order:
        group = phase_groups[key]
        full_body.append(('phase', group['name']))
        for d in group['items']:
            full_body.append(('item', d))

    # --- Paginate: max ITEMS_PER_PAGE deliverable rows per slide ---
    # Phase headers carry over to the next page when a group spans the boundary.
    pages = []          # list of body-row lists, one per slide
    current_page = []
    page_item_count = 0
    last_phase_seen = None  # most recent phase header for carry-over

    for entry in full_body:
        kind, payload = entry
        if kind == 'phase':
            last_phase_seen = entry
            current_page.append(entry)
        else:
            # Deliverable item — flush page when limit reached
            if page_item_count >= ITEMS_PER_PAGE:
                pages.append(current_page)
                current_page = []
                page_item_count = 0
                # Carry the current phase header so context is clear on the new slide
                if last_phase_seen is not None:
                    current_page.append(last_phase_seen)
            current_page.append(entry)
            page_item_count += 1

    if current_page:
        pages.append(current_page)

    num_pages = len(pages)

    # --- Render one slide per page ---
    first_slide = None
    for page_idx, body in enumerate(pages):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        if first_slide is None:
            first_slide = slide
        add_accent_bar(slide, primary_color, top=0)

        # Title
        page_label = f"  (Page {page_idx + 1} of {num_pages})" if num_pages > 1 else ""
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(10), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"Deliverables Tracker{page_label}"
        set_font(run, size=22, bold=True, color=primary_color)

        # Summary stats (overall counts, repeated on every slide for quick reference)
        summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(12), Inches(0.35))
        stf = summary_box.text_frame
        sp = stf.paragraphs[0]
        srun = sp.add_run()
        srun.text = summary_text
        set_font(srun, size=max(10, MIN_FONT), bold=False, color='#666666')

        table_top = Inches(1.4)
        rows = len(body) + 1  # +1 for column header row

        table_shape = slide.shapes.add_table(
            rows, cols,
            Inches(0.3), table_top,
            sum(col_widths), Inches(0.34 * rows)
        )
        table = table_shape.table
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

        # Column headers
        headers = ['Deliverable', 'Owner', 'Status', 'Target Date', 'Delivered Date']
        for i, h in enumerate(headers):
            set_cell_text(table.cell(0, i), h, size=max(9, MIN_FONT), bold=True,
                          color='#FFFFFF', bg_color=primary_color, alignment=PP_ALIGN.LEFT)

        alt = 0
        for idx, (kind, payload) in enumerate(body):
            r = idx + 1
            if kind == 'phase':
                merged = table.cell(r, 0)
                merged.merge(table.cell(r, cols - 1))
                phase_name = payload
                if len(phase_name) > 70:
                    phase_name = phase_name[:67] + '...'
                set_cell_text(merged, phase_name, size=max(9, MIN_FONT), bold=True,
                              color=primary_color, bg_color='#ECE6F8', alignment=PP_ALIGN.LEFT)
                alt = 0
                continue

            d = payload
            alt += 1
            bg = '#FFFFFF' if alt % 2 == 1 else '#F8F8FC'

            name = d.get('name', '')
            if len(name) > 48:
                name = name[:45] + '...'
            set_cell_text(table.cell(r, 0), name, size=max(MIN_FONT, MIN_FONT), bg_color=bg)
            set_cell_text(table.cell(r, 1), d.get('ownerName', ''), size=MIN_FONT, bg_color=bg)

            status = d.get('status', 'not-started')
            status_label = status.replace('-', ' ').title()
            s_color = status_colors.get(status, '#9CA3AF')
            set_cell_text(table.cell(r, 2), status_label, size=MIN_FONT, bold=True,
                          color=s_color, bg_color=bg)

            set_cell_text(table.cell(r, 3), d.get('targetDate', ''), size=MIN_FONT, bg_color=bg)
            set_cell_text(table.cell(r, 4), d.get('deliveredDate', ''), size=MIN_FONT, bg_color=bg)

        # Page footer
        if num_pages > 1:
            footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.3))
            ftf = footer_box.text_frame
            fp = ftf.paragraphs[0]
            frun = fp.add_run()
            items_on_page = sum(1 for kind, _ in body if kind == 'item')
            frun.text = f"Page {page_idx + 1} of {num_pages}  ·  {items_on_page} deliverables shown  ·  {total} total"
            set_font(frun, size=MIN_FONT, italic=True, color='#999999')

    return first_slide


def _remap_rids_in_element(element, rId_map):
    """Walk an lxml element tree and replace rId attribute values according to rId_map."""
    if not rId_map:
        return
    for elem in element.iter():
        for attr_name, attr_val in list(elem.attrib.items()):
            if attr_val in rId_map:
                elem.set(attr_name, rId_map[attr_val])


def _normalize_fonts_in_element(element, explicit_font=FONT_NAME):
    """
    Ensure all text in a copied slide element uses the project font (Avenir Next LT Pro).

    Two cases are handled:
    1. Explicit theme-font tokens (+mj-lt, +mn-lt, etc.) on <a:latin>/<a:ea>/<a:cs>
       elements — these resolve against the destination theme (Calibri by default)
       instead of the source theme.  The token is replaced with an explicit name.
    2. Run-property elements (<a:rPr>, <a:defRPr>, <a:endParaRPr>) with NO <a:latin>
       child — font is fully inherited through the theme chain, again landing on
       Calibri.  An explicit <a:latin> element is injected to break the inheritance.
    """
    from lxml import etree as _etree
    A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    LATIN_TAG   = f'{{{A}}}latin'
    FONT_TAGS   = {LATIN_TAG, f'{{{A}}}ea', f'{{{A}}}cs'}
    RPR_TAGS    = {f'{{{A}}}rPr', f'{{{A}}}defRPr', f'{{{A}}}endParaRPr'}

    for elem in element.iter():
        if elem.tag in FONT_TAGS:
            typeface = elem.get('typeface', '')
            if typeface.startswith('+') or typeface == '':
                elem.set('typeface', explicit_font)
        elif elem.tag in RPR_TAGS:
            if elem.find(LATIN_TAG) is None:
                latin_el = _etree.SubElement(elem, LATIN_TAG)
                latin_el.set('typeface', explicit_font)


def _slide_is_layout_shell(slide_element, P, A):
    """
    Return True when the slide's <p:cSld> contains nothing but empty placeholder
    shapes with no explicit position and no text content.

    PowerPoint templates created from a custom slide layout export this way: the
    entire visual design lives in the layout, and the slide itself is just a thin
    shell of empty <p:ph> wrappers.
    """
    cSld = slide_element.find(f'{{{P}}}cSld')
    if cSld is None:
        return True
    spTree = cSld.find(f'{{{P}}}spTree')
    if spTree is None:
        return True
    for child in spTree:
        local = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if local in ('nvGrpSpPr', 'grpSpPr', 'extLst'):
            continue
        if local != 'sp':
            return False  # pic/graphicFrame/etc. → has real content
        nvSpPr = child.find(f'{{{P}}}nvSpPr')
        nvPr   = nvSpPr.find(f'{{{P}}}nvPr') if nvSpPr is not None else None
        ph_el  = nvPr.find(f'{{{P}}}ph')     if nvPr   is not None else None
        if ph_el is None:
            return False  # non-placeholder shape → has real content
        spPr = child.find(f'{{{P}}}spPr')
        if spPr is not None and spPr.find(f'{{{A}}}xfrm') is not None:
            return False  # placeholder with explicit position → real content
        txBody = child.find(f'{{{P}}}txBody')
        if txBody is not None:
            for t_el in txBody.iter(f'{{{A}}}t'):
                if t_el.text and t_el.text.strip():
                    return False  # placeholder with typed text → real content
    return True


def _inject_placeholder_text(cSld_el, inject_texts, P, A):
    """
    Replace text inside specific placeholder shapes with actual report content.

    Must be called BEFORE _strip_placeholder_tags so <p:ph> is still present
    and can be used to identify which shape is the title vs the subtitle.

    inject_texts keys:
        'title'    → shape with <p:ph type="title">
        'subtitle' → shape with <p:ph idx="1"> (body / subtitle placeholder)
        'number'   → any shape (placeholder OR free-form text box) whose sole
                     text content is exactly "01" — replaces it with the supplied
                     value (e.g. "02", "03").  This handles templates where the
                     section-number badge is a hardcoded text box, not a placeholder.
    """
    import re as _re
    from lxml import etree as _etree

    if not inject_texts:
        return
    spTree = cSld_el.find(f'{{{P}}}spTree')
    if spTree is None:
        return

    number_value = inject_texts.get('number')
    _num_pat = _re.compile(r'^\s*0\d\s*$')  # matches "01", "02", … "09"

    for sp in spTree.findall(f'{{{P}}}sp'):
        nvSpPr = sp.find(f'{{{P}}}nvSpPr')
        if nvSpPr is None:
            continue

        txBody = sp.find(f'{{{P}}}txBody')
        if txBody is None:
            continue

        # ── 'number' injection: scan every shape for a run that is exactly "0N" ──
        if number_value:
            all_runs = txBody.findall(f'.//{{{A}}}r')
            all_texts = [r.find(f'{{{A}}}t') for r in all_runs]
            full_text = ''.join(t.text or '' for t in all_texts if t is not None)
            if _num_pat.match(full_text):
                for t_el in all_texts:
                    if t_el is not None:
                        t_el.text = number_value if t_el == all_texts[0] else ''
                continue  # done with this shape; skip placeholder logic below

        # ── title / subtitle injection (placeholder shapes only) ──
        nvPr = nvSpPr.find(f'{{{P}}}nvPr')
        if nvPr is None:
            continue
        ph_el = nvPr.find(f'{{{P}}}ph')
        if ph_el is None:
            continue

        ph_type = ph_el.get('type', '')
        ph_idx  = ph_el.get('idx',  '0')

        if ph_type == 'title':
            new_text = inject_texts.get('title')
        elif ph_idx == '1':
            new_text = inject_texts.get('subtitle')
        else:
            continue

        if new_text is None:
            continue

        # Preserve the run properties from the first run (bold, colour, size)
        first_rPr = None
        first_para = txBody.find(f'{{{A}}}p')
        if first_para is not None:
            first_run = first_para.find(f'{{{A}}}r')
            if first_run is not None:
                first_rPr = first_run.find(f'{{{A}}}rPr')

        # Support either a plain string or a list of strings (one paragraph each)
        lines = [new_text] if isinstance(new_text, str) else [l for l in new_text if l]
        if not lines:
            lines = [' ']  # keep a non-empty run so the placeholder box stays editable

        # Remove all existing paragraphs and replace with one paragraph per line
        for p_el in txBody.findall(f'{{{A}}}p'):
            txBody.remove(p_el)
        for line in lines:
            new_p = _etree.SubElement(txBody, f'{{{A}}}p')
            new_r = _etree.SubElement(new_p, f'{{{A}}}r')
            if first_rPr is not None:
                new_r.insert(0, copy.deepcopy(first_rPr))
            new_t = _etree.SubElement(new_r, f'{{{A}}}t')
            new_t.text = line


def _strip_placeholder_tags(cSld_el, P):
    """
    Remove <p:ph> from every shape in cSld_el's spTree.

    When layout content is promoted into a slide, the <p:ph> elements cause
    PowerPoint to merge the shape's styling with the *destination* layout's
    placeholder — which is the wrong layout and provides wrong positions/colours.
    Stripping <p:ph> makes every shape a standalone freeform with its own
    explicit attributes and no external inheritance.
    """
    spTree = cSld_el.find(f'{{{P}}}spTree')
    if spTree is None:
        return
    for sp in list(spTree):
        nvSpPr = sp.find(f'{{{P}}}nvSpPr')
        if nvSpPr is None:
            continue
        nvPr = nvSpPr.find(f'{{{P}}}nvPr')
        if nvPr is None:
            continue
        ph_el = nvPr.find(f'{{{P}}}ph')
        if ph_el is not None:
            nvPr.remove(ph_el)


def copy_first_slide(src_pptx_path, dest_prs, _part_cache=None, inject_texts=None):
    """
    Copy the first slide from src_pptx_path into dest_prs, fully flattening the
    source master/layout chain so the destination slide is visually complete
    without needing the source master.

    Many branded template PPTX files are created from a custom slide layout where
    the entire visual design (background, positioned shapes, logo, text content)
    lives in the *layout*, not the slide itself.  The slide XML is just a thin
    shell of empty placeholder wrappers.  This function detects that pattern and
    uses the layout's <p:cSld> as the content source instead of the empty slide.

    Strategy:
    1. Detect whether the slide is an empty layout shell.
    2. Use the layout's <p:cSld> as content when the slide is empty; otherwise
       use the slide's own <p:cSld> with a background fallback from layout/master.
    3. Strip <p:ph> tags from copied shapes so the destination layout cannot
       override the source positions or styling.
    4. Copy image parts from slide + layout + master.
    5. Remap rIds, normalise fonts (replace theme tokens), install on blank slide.

    _part_cache: optional dict shared across calls with the same template file so
    images are not duplicated (used for repeated section-header slides).

    Returns True on success, False on failure.
    """
    try:
        from pptx.parts.image import ImagePart
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT

        P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        src_prs = Presentation(src_pptx_path)
        if not src_prs.slides:
            print(f"[PPTX_TEMPLATE] Source PPTX has no slides: {src_pptx_path}", file=sys.stderr)
            return False

        src_slide  = src_prs.slides[0]
        src_layout = src_slide.slide_layout
        src_master = src_layout.slide_master

        dest_pkg = dest_prs.part.package
        cache = _part_cache if _part_cache is not None else {}

        # ── Step 1: choose content source ─────────────────────────────────────
        # When the slide is an empty layout shell (all placeholders, no text,
        # no explicit positions) the real design lives in the slide layout.
        # Use the layout's <p:cSld> in that case; otherwise use the slide's own.
        is_shell = _slide_is_layout_shell(src_slide.element, P, A)
        if is_shell:
            content_el = src_layout.element
            print(f"[PPTX_TEMPLATE] Slide is empty layout shell — using layout content", file=sys.stderr)
        else:
            content_el = src_slide.element

        src_cSld = content_el.find(f'{{{P}}}cSld')
        if src_cSld is None:
            print(f"[PPTX_TEMPLATE] No <p:cSld> in content source: {src_pptx_path}", file=sys.stderr)
            return False
        new_cSld = copy.deepcopy(src_cSld)

        # Inject actual report content into placeholder shapes BEFORE stripping
        # the <p:ph> tags — the ph type/idx is how we identify title vs subtitle.
        if inject_texts:
            _inject_placeholder_text(new_cSld, inject_texts, P, A)

        # Strip <p:ph> placeholders so destination layout cannot override positions.
        _strip_placeholder_tags(new_cSld, P)

        # ── Step 2: if there is still no <p:bg>, inherit from layout/master ───
        # (only relevant when using slide content that lacks an explicit background)
        BG_TAG = f'{{{P}}}bg'
        if new_cSld.find(BG_TAG) is None:
            inherited_bg = None
            for ancestor_el in [src_layout.element, src_master.element]:
                ancestor_cSld = ancestor_el.find(f'{{{P}}}cSld')
                if ancestor_cSld is not None:
                    candidate = ancestor_cSld.find(BG_TAG)
                    if candidate is not None:
                        inherited_bg = copy.deepcopy(candidate)
                        break
            if inherited_bg is not None:
                spTree = new_cSld.find(f'{{{P}}}spTree')
                insert_at = list(new_cSld).index(spTree) if spTree is not None else 0
                new_cSld.insert(insert_at, inherited_bg)
                print(f"[PPTX_TEMPLATE] Inherited <p:bg> from layout/master into slide", file=sys.stderr)

        # ── Step 3: add blank slide to destination ─────────────────────────────
        blank_layout = (
            dest_prs.slide_layouts[6]
            if len(dest_prs.slide_layouts) > 6
            else dest_prs.slide_layouts[-1]
        )
        new_slide = dest_prs.slides.add_slide(blank_layout)
        dest_part = new_slide.part

        # ── Step 4: copy image parts from slide + layout + master ──────────────
        # Branded templates often embed background images only in the layout or
        # master; collecting from all three levels ensures nothing is missed.
        rId_map = {}

        def _copy_images_from(part):
            for old_rId, rel in part.rels.items():
                try:
                    if rel.is_external:
                        continue
                    if not isinstance(rel.target_part, ImagePart):
                        continue
                    cache_key = str(rel.target_part.partname)
                    if cache_key in cache:
                        img_part = cache[cache_key]
                    else:
                        img_part = ImagePart.new(dest_pkg, rel.target_part.image)
                        cache[cache_key] = img_part
                    new_rId = dest_part.relate_to(img_part, rel.reltype)
                    if old_rId != new_rId:
                        rId_map[old_rId] = new_rId
                except Exception as rel_err:
                    print(f"[PPTX_TEMPLATE] Image rel {old_rId} skipped: {rel_err}", file=sys.stderr)

        slide_img_count  = sum(1 for r in src_slide.part.rels.values() if not r.is_external and isinstance(r.target_part, ImagePart))
        layout_img_count = sum(1 for r in src_layout.part.rels.values() if not r.is_external and isinstance(r.target_part, ImagePart))
        master_img_count = sum(1 for r in src_master.part.rels.values() if not r.is_external and isinstance(r.target_part, ImagePart))

        _copy_images_from(src_slide.part)
        _copy_images_from(src_layout.part)
        _copy_images_from(src_master.part)

        # ── Step 5: remap rIds, fix theme fonts, then install the new <p:cSld> ──
        _remap_rids_in_element(new_cSld, rId_map)
        # Theme-font tokens (e.g. +mj-lt, +mn-lt) resolve against the destination
        # presentation's theme, which is the default Office theme (Calibri).  Replace
        # them with the explicit project font so headings use Avenir Next LT Pro.
        _normalize_fonts_in_element(new_cSld)

        dest_sld = dest_part._element
        dst_cSld = dest_sld.find(f'{{{P}}}cSld')
        if dst_cSld is not None:
            dest_sld.replace(dst_cSld, new_cSld)
        else:
            dest_sld.insert(0, new_cSld)

        # Preserve the source's colour-map override when present — it controls
        # how theme accent colours are inherited from the slide master.
        src_clrMapOvr = src_slide.element.find(f'{{{P}}}clrMapOvr')
        if src_clrMapOvr is not None:
            dst_clrMapOvr = dest_sld.find(f'{{{P}}}clrMapOvr')
            new_clrMapOvr = copy.deepcopy(src_clrMapOvr)
            if dst_clrMapOvr is not None:
                dest_sld.replace(dst_clrMapOvr, new_clrMapOvr)
            else:
                dest_sld.append(new_clrMapOvr)

        print(
            f"[PPTX_TEMPLATE] Copied slide from {os.path.basename(src_pptx_path)} "
            f"(imgs: slide={slide_img_count}, layout={layout_img_count}, "
            f"master={master_img_count}; {len(rId_map)} rIds remapped)",
            file=sys.stderr,
        )
        return True

    except Exception as e:
        import traceback
        print(f"[PPTX_TEMPLATE] Failed to copy slide from {src_pptx_path}: {e}", file=sys.stderr)
        print(f"[PPTX_TEMPLATE] Traceback: {traceback.format_exc()}", file=sys.stderr)
        return False


def create_exec_title_slide(prs, data, primary_color, secondary_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = hex_to_rgb(primary_color)
    bg_shape.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.0), SLIDE_WIDTH, Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = hex_to_rgb(secondary_color)
    accent.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "Executive Narrative"
    set_font(run, size=36, bold=True, color=RGBColor(255, 255, 255))
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = "PRACTICE SUMMARY"
    set_font(run2, size=20, color=RGBColor(220, 220, 220))

    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1.2))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    run = p.add_run()
    tenant_name = data.get('tenantName', '')
    run.text = tenant_name
    set_font(run, size=16, bold=True, color=RGBColor(255, 255, 255))

    period_start = data.get('periodStart', '')
    period_end = data.get('periodEnd', '')
    report_date = data.get('reportDate', datetime.now().strftime('%B %d, %Y'))
    period_text = f"Period: {period_start} to {period_end}" if period_start and period_end else report_date
    p2 = tf2.add_paragraph()
    run2 = p2.add_run()
    run2.text = period_text
    set_font(run2, size=14, color=RGBColor(230, 230, 230))

    p3 = tf2.add_paragraph()
    run3 = p3.add_run()
    run3.text = report_date
    set_font(run3, size=12, color=RGBColor(200, 200, 200))

    logo_path = data.get('logoPath')
    if logo_path and os.path.exists(logo_path):
        try:
            slide.shapes.add_picture(logo_path, Inches(10.5), Inches(4.0), height=Inches(0.8))
        except Exception:
            pass
    return slide


def create_exec_financial_slide(prs, data, primary_color, secondary_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide, primary_color, top=0)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(10), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Financial Performance"
    set_font(run, size=24, bold=True, color=primary_color)

    stats = data.get('stats', {})
    metrics = [
        ('Billable Hours', f"{stats.get('billableHours', 0):.1f}"),
        ('Total Hours', f"{stats.get('totalHours', 0):.1f}"),
        ('Revenue', f"${stats.get('totalRevenue', 0):,.0f}"),
        ('Expenses', f"${stats.get('totalExpenses', 0):,.0f}"),
        ('Active Projects', str(stats.get('activeProjects', 0))),
        ('Estimates Created', str(stats.get('estimatesCreated', 0))),
        ('Milestones Completed', str(stats.get('milestonesCompleted', 0))),
        ('Status Reports', str(stats.get('statusReportsPublished', 0))),
    ]

    col_count = 4
    row_count = 2
    card_w = Inches(2.7)
    card_h = Inches(1.3)
    start_x = Inches(0.8)
    start_y = Inches(1.2)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)

    for i, (label, value) in enumerate(metrics):
        col = i % col_count
        row = i // col_count
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(245, 245, 245)
        card.line.color.rgb = RGBColor(220, 220, 220)
        card.line.width = Pt(0.5)

        val_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), card_w - Inches(0.4), Inches(0.6))
        vf = val_box.text_frame
        vp = vf.paragraphs[0]
        vr = vp.add_run()
        vr.text = value
        set_font(vr, size=22, bold=True, color=primary_color)

        lbl_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.8), card_w - Inches(0.4), Inches(0.4))
        lf = lbl_box.text_frame
        lp = lf.paragraphs[0]
        lr = lp.add_run()
        lr.text = label
        set_font(lr, size=10, color='#666666')

    utilization = stats.get('totalHours', 0)
    if utilization > 0:
        pct = (stats.get('billableHours', 0) / utilization) * 100
        util_text = f"Utilization Rate: {pct:.0f}%"
    else:
        util_text = "Utilization Rate: N/A"

    revenue = stats.get('totalRevenue', 0)
    if revenue > 0:
        margin = ((revenue - stats.get('totalExpenses', 0)) / revenue) * 100
        margin_text = f"Gross Margin: {margin:.1f}%"
    else:
        margin_text = "Gross Margin: N/A"

    summary_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(10), Inches(0.5))
    sf = summary_box.text_frame
    sp = sf.paragraphs[0]
    sr = sp.add_run()
    sr.text = f"{util_text}   |   {margin_text}   |   Active Assignments: {stats.get('activeAssignments', 0)}"
    set_font(sr, size=11, bold=True, color='#444444')

    return slide


def create_exec_narrative_slides(prs, data, sections, primary_color, secondary_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide, primary_color, top=0)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(10), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Practice Overview"
    set_font(run, size=24, bold=True, color=primary_color)

    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.5), Inches(6.0))
    tf = content_box.text_frame
    tf.word_wrap = True

    found = False
    for key in ['Practice Overview', 'Progress Summary', 'Executive Summary', 'Summary', 'Overview']:
        if key in sections and sections[key].strip():
            render_markdown_text(tf, sections[key], primary_color, size=11, start_fresh=True)
            found = True
            break

    if not found:
        all_text = '\n'.join(sections.values())
        if all_text.strip():
            lines = all_text.strip().split('\n')
            overview_lines = lines[:40]
            render_markdown_text(tf, '\n'.join(overview_lines), primary_color, size=11, start_fresh=True)

    return slide


def create_exec_raidd_summary_slide(prs, data, sections, primary_color, secondary_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide, primary_color, top=0)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(10), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Risks, Issues & Key Decisions"
    set_font(run, size=24, bold=True, color=primary_color)

    stats = data.get('stats', {})
    counts_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(10), Inches(0.4))
    cf = counts_box.text_frame
    cp = cf.paragraphs[0]
    cr = cp.add_run()
    cr.text = f"Open Risks: {stats.get('openRisks', 0)}   |   Open Issues: {stats.get('openIssues', 0)}   |   Open Actions: {stats.get('openActions', 0)}"
    set_font(cr, size=11, bold=True, color='#444444')

    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for key in ['Risks, Issues & Key Decisions (RAIDD)', 'RAIDD', 'Risk Summary', 'Risks and Issues']:
        if key in sections and sections[key].strip():
            render_markdown_text(tf, sections[key], primary_color, size=10, start_fresh=True)
            return slide

    raidd_items = data.get('raiddHighPriority', [])
    if raidd_items:
        first = True
        for item in raidd_items[:12]:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(4)
            run = p.add_run()
            prefix = f"[{item.get('type', '').upper()}] "
            ref = item.get('refNumber', '')
            if ref:
                prefix += f"{ref} "
            run.text = f"• {prefix}{item.get('title', '')}"
            set_font(run, size=10, bold=True, color=primary_color)

            detail = item.get('impact') or item.get('description') or ''
            if detail:
                p2 = tf.add_paragraph()
                p2.space_before = Pt(1)
                r2 = p2.add_run()
                r2.text = f"  {detail}"
                set_font(r2, size=9, color='#444444')
    else:
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = "No high-priority risks or issues in this period."
        set_font(r, size=11, color='#666666')

    return slide


def create_exec_outlook_slide(prs, data, sections, primary_color, secondary_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide, primary_color, top=0)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(10), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Outlook & Recommendations"
    set_font(run, size=24, bold=True, color=primary_color)

    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.5), Inches(6.0))
    tf = content_box.text_frame
    tf.word_wrap = True

    for key in ['Upcoming Activities', 'Outlook', 'Next Steps', 'Recommendations', 'Looking Ahead', 'Outlook & Recommendations']:
        if key in sections and sections[key].strip():
            render_markdown_text(tf, sections[key], primary_color, size=11, start_fresh=True)
            return slide

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "See narrative for forward-looking commentary."
    set_font(r, size=11, color='#666666')
    return slide


def generate_executive_narrative_pptx(data, output_path):
    primary_color = data.get('primaryColor', '#810FFB')
    secondary_color = data.get('secondaryColor', '#E60CB3')

    narrative = data.get('narrative', '')
    print(f"[EXEC-PPTX] narrative length: {len(narrative)}, first 300: {narrative[:300]}", file=sys.stderr)

    EXEC_SECTION_ALIASES = {
        'Practice Overview': [
            'practice overview', 'executive summary', 'overview', 'summary',
            'progress summary', 'period overview', 'portfolio overview',
        ],
        'Financial Performance': [
            'financial performance', 'financial summary', 'financials',
            'financial highlights', 'revenue & utilization', 'revenue and utilization',
        ],
        'Key Accomplishments': [
            'key accomplishments', 'accomplishments', 'highlights',
            'key highlights', 'achievements', 'progress & accomplishments',
        ],
        'Risks, Issues & Key Decisions (RAIDD)': [
            'risks, issues & key decisions', 'raidd', 'risks and issues',
            'risks, issues & key decisions (raidd)', 'risk summary',
            'raidd summary', 'raidd log', 'risk & issue summary',
        ],
        'Upcoming Activities': [
            'upcoming activities', 'next steps', 'outlook', 'recommendations',
            'looking ahead', 'outlook & recommendations', 'forward look',
        ],
    }

    def _normalize_exec_section(raw_name):
        lower = raw_name.lower().strip()
        lower = re.sub(r'\s*\(.*?\)\s*$', '', lower).strip()
        for canonical, aliases in EXEC_SECTION_ALIASES.items():
            if lower in aliases:
                return canonical
            for alias in aliases:
                if alias in lower or lower in alias:
                    return canonical
        return raw_name

    sections = {}
    if narrative:
        current_section = None
        current_content = []
        for line in narrative.split('\n'):
            header_match = re.match(r'^(#{1,3})\s+(.+)$', line)
            if header_match:
                if current_section:
                    sections[current_section] = '\n'.join(current_content).strip()
                raw_name = header_match.group(2).strip()
                current_section = _normalize_exec_section(raw_name)
                current_content = []
            else:
                current_content.append(line)
        if current_section:
            sections[current_section] = '\n'.join(current_content).strip()

    print(f"[EXEC-PPTX] Sections found: {list(sections.keys())}", file=sys.stderr)

    title_template_path = data.get('titleTemplatePath')
    section_template_path = data.get('sectionTemplatePath')
    closing_template_path = data.get('closingTemplatePath')

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    section_cache = {}

    period_start = data.get('periodStart', '')
    period_end = data.get('periodEnd', '')
    report_date = data.get('reportDate', datetime.now().strftime('%B %d, %Y'))
    period_label = (f"Period: {period_start} to {period_end}"
                    if period_start and period_end else report_date)
    tenant_name = data.get('tenantName', '')
    subtitle_lines = [l for l in [tenant_name, period_label] if l]
    title_inject = {
        'title': 'Executive Narrative',
        'subtitle': subtitle_lines,
    }

    if title_template_path and os.path.exists(title_template_path):
        print(f"[EXEC-PPTX] Using title template: {title_template_path}", file=sys.stderr)
        if not copy_first_slide(title_template_path, prs, inject_texts=title_inject):
            create_exec_title_slide(prs, data, primary_color, secondary_color)
    else:
        create_exec_title_slide(prs, data, primary_color, secondary_color)

    _exec_section_counter = [0]

    def insert_section_header(section_name):
        _exec_section_counter[0] += 1
        if section_template_path and os.path.exists(section_template_path):
            ok = copy_first_slide(
                section_template_path, prs,
                _part_cache=section_cache,
                inject_texts={
                    'title':    section_name,
                    'number':   f"{_exec_section_counter[0]:02d}",
                    'subtitle': ' ',
                },
            )
            if not ok:
                print(f"[EXEC-PPTX] Section header failed for '{section_name}'", file=sys.stderr)

    insert_section_header('Financial Performance')
    create_exec_financial_slide(prs, data, primary_color, secondary_color)

    insert_section_header('Practice Overview')
    create_exec_narrative_slides(prs, data, sections, primary_color, secondary_color)

    insert_section_header('Risks & Issues')
    create_exec_raidd_summary_slide(prs, data, sections, primary_color, secondary_color)

    insert_section_header('Outlook')
    create_exec_outlook_slide(prs, data, sections, primary_color, secondary_color)

    if closing_template_path and os.path.exists(closing_template_path):
        print(f"[EXEC-PPTX] Using closing template: {closing_template_path}", file=sys.stderr)
        ok = copy_first_slide(closing_template_path, prs)
        if not ok:
            print(f"[EXEC-PPTX] Closing template failed", file=sys.stderr)

    prs.save(output_path)
    return output_path


def generate_pptx(data, output_path):
    primary_color = data.get('primaryColor', '#810FFB')
    secondary_color = data.get('secondaryColor', '#E60CB3')

    ai_report = data.get('aiReport', '')
    print(f"[PPTX] aiReport length: {len(ai_report)}, first 300: {ai_report[:300]}", file=sys.stderr)
    sections = parse_markdown_sections(ai_report) if ai_report else {}
    print(f"[PPTX] Sections found: {list(sections.keys())}, content lengths: {({k: len(v) for k, v in sections.items()})}", file=sys.stderr)

    title_template_path = data.get('titleTemplatePath')
    section_template_path = data.get('sectionTemplatePath')
    closing_template_path = data.get('closingTemplatePath')

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    section_cache = {}

    # Build title slide text from report data (mirrors what create_title_slide shows)
    period_start = data.get('periodStart', '')
    period_end   = data.get('periodEnd',   '')
    report_date  = data.get('reportDate',  datetime.now().strftime('%B %d, %Y'))
    period_label = (f"Period: {period_start} to {period_end}"
                    if period_start and period_end else report_date)
    client_name  = data.get('clientName', '')
    pm_name      = data.get('pmName', '')
    pm_line      = f"Project Manager: {pm_name}" if pm_name else ''
    subtitle_lines = [l for l in [client_name, period_label, pm_line] if l]
    title_inject = {
        'title':    data.get('projectName', 'Status Report'),
        'subtitle': subtitle_lines,
    }

    # Title slide: use template if provided, otherwise generate programmatically
    if title_template_path and os.path.exists(title_template_path):
        print(f"[PPTX_TEMPLATE] Using title template: {title_template_path}", file=sys.stderr)
        if not copy_first_slide(title_template_path, prs, inject_texts=title_inject):
            create_title_slide(prs, data, primary_color, secondary_color)
    else:
        create_title_slide(prs, data, primary_color, secondary_color)

    _section_counter = [0]

    def insert_section_header(section_name):
        _section_counter[0] += 1
        if section_template_path and os.path.exists(section_template_path):
            ok = copy_first_slide(
                section_template_path, prs,
                _part_cache=section_cache,
                inject_texts={
                    'title':    section_name,
                    'number':   f"{_section_counter[0]:02d}",
                    'subtitle': ' ',
                },
            )
            if not ok:
                print(f"[PPTX_TEMPLATE] Section header insertion failed before '{section_name}' — continuing without template slide", file=sys.stderr)

    # Section header before Progress Summary
    insert_section_header('Progress Summary')
    create_progress_summary_slide(prs, data, sections, primary_color, secondary_color)
    create_accomplishments_slide(prs, data, sections, primary_color, secondary_color)

    # Section header before RAIDD
    insert_section_header('RAIDD')
    create_raidd_slides(prs, data, sections, primary_color, secondary_color)

    # Section header before Upcoming
    insert_section_header('Upcoming Activities')
    create_upcoming_slide(prs, data, sections, primary_color, secondary_color)
    create_deliverables_slide(prs, data, primary_color, secondary_color)
    create_timeline_slide(prs, data, primary_color, secondary_color)

    # Section header before Project Plan
    insert_section_header('Project Plan')
    create_project_plan_slides(prs, data, primary_color, secondary_color)

    # Closing slide: append after last slide
    if closing_template_path and os.path.exists(closing_template_path):
        print(f"[PPTX_TEMPLATE] Using closing template: {closing_template_path}", file=sys.stderr)
        ok = copy_first_slide(closing_template_path, prs)
        if not ok:
            print(f"[PPTX_TEMPLATE] Closing template insertion failed — deck will end without closing slide", file=sys.stderr)

    prs.save(output_path)
    return output_path

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage: generate_status_report_pptx.py <output_path> [--executive-narrative]", file=sys.stderr)
        sys.exit(1)

    output_path = sys.argv[1]
    is_exec_narrative = '--executive-narrative' in sys.argv
    input_data = json.load(sys.stdin)

    if is_exec_narrative:
        result = generate_executive_narrative_pptx(input_data, output_path)
    else:
        result = generate_pptx(input_data, output_path)
    print(json.dumps({"success": True, "path": result}))
